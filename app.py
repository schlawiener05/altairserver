from flask import Flask
from flask import jsonify, send_file
from flask import render_template
import altair as alt
from vega_datasets import data

from io import StringIO
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)

cars = alt.load_dataset('cars')
source = data.barley()

@app.route('/')
def hello_world():
    return 'Flask Dockerized'

@app.route('/index')
def index():
    user = {'username': 'Miguel'}
    achart = create_chart()
    return render_template('index.html', title='Home', user=user, altairchart = achart)

@app.route('/test')
def test():
    cjson = create_chart()
    return jsonify(cjson)

@app.route('/testdata')
def testdata():
    return jsonify(source.to_dict('records'))

@app.route('/testppt')
def testppt():
    out_file = BytesIO()

    imagefile = BytesIO()
    # imagefile = StringIO()
    c = create_chart_object()
    c.save(imagefile, 'png',scale_factor=1.9)
    imagefile.seek(0)

    path = 'template.pptx'
    prs = Presentation(path)
    # blank_slide_layout = prs.slide_layouts[6]
    # slide = prs.slides.add_slide(blank_slide_layout)
    title_only_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = 'Adding an AutoShape'

    left = Pt(50)
    top = Pt(300)
    width = Pt(700)
    height = Pt(350) 
    txBox = shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    tf.word_wrap = True
    tf.clear()  # remove any existing paragraphs, leaving one empty one
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = "LS 1234567"
    p.font.bold = True
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.text = "This is a second paragraph containing some comment"
    # top = Inches(1.5)
    # left = Inches(5)
    top = Pt(300)
    left = Pt(900)
    pic = shapes.add_picture(imagefile, left, top, height = Pt(500))
    
    symptoms = [{"abbr":"A", "symp":"bubu"},{"abbr":"B", "symp":"bibi"},{"abbr":"CA", "symp":"lala"}]
    rows = len(symptoms)+1
    cols = 2
    left = Pt(50)
    top = Pt(650)
    width = Pt(350)
    height = Pt(150)

    table = shapes.add_table(rows, cols, left, top, width, height).table
    # set column widths
    table.columns[0].width = Pt(70)
    table.columns[1].width = Pt(300)

    table.cell(0, 0).text = 'Abbreviation'
    table.cell(0, 1).text = 'Symptom'

    for i in range(rows-1):
        table.cell(i, 0).text = symptoms[i]['abbr']
        table.cell(i, 1).text = symptoms[i]['symp']

    prs.save(out_file)
    out_file.seek(0)
    return send_file(out_file, attachment_filename="testing.pptx", as_attachment=True)

def create_chart():
    c=alt.Chart(source).mark_bar().encode(
        x='yield:Q',
        y='year:O',
        color='year:N',
        row='site:N',
        tooltip=['yield', 'year']
    )
    cjson = c.to_dict()
    return cjson

def create_chart_object():
    c=alt.Chart(source).mark_bar().encode(
        x='yield:Q',
        y='year:O',
        color='year:N',
        row='site:N',
        tooltip=['yield', 'year']
    )
    return c

if __name__ == '__main__':
    app.run(debug=True,host='0.0.0.0')